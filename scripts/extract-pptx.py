#!/usr/bin/env python3
"""
extract-pptx.py — Extract all content from a PowerPoint (.pptx) file.
Outputs a JSON structure with slides, text, images, and speaker notes.

Usage:
    python extract-pptx.py <input.pptx> [output_dir]

Requires python-pptx (auto-installed if missing).
"""

import json
import os
import subprocess
import sys


# ─── Auto-install python-pptx if missing (BUG-FIX 19) ───────────────────────

def _ensure_pptx() -> None:
    try:
        import pptx  # noqa: F401
    except ImportError:
        print("python-pptx not found — installing...")
        try:
            subprocess.run(
                [sys.executable, "-m", "pip", "install",
                 "--break-system-packages", "-q", "python-pptx"],
                check=True, timeout=120,
            )
            print("python-pptx installed.\n")
        except Exception as e:
            print(f"Auto-install failed: {e}")
            print("Run manually:  pip install python-pptx")
            sys.exit(1)


# ─── Helper: recurse into group shapes (BUG-FIX 16) ─────────────────────────

def _iter_shapes(shapes):
    """Yield every shape, recursing into groups."""
    from pptx.util import Pt  # noqa: F401 — import used downstream
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)


# ─── Core extraction ─────────────────────────────────────────────────────────

def extract_pptx(file_path: str, output_dir: str = ".") -> list:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    # BUG-FIX 18: catch corrupt / wrong-type files with a clear message
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"Error opening '{file_path}': {e}")
        print("Make sure the file is a valid .pptx (not .ppt or a renamed PDF).")
        sys.exit(1)

    assets_dir = os.path.join(output_dir, "assets")
    os.makedirs(assets_dir, exist_ok=True)

    slides_data = []

    for slide_num, slide in enumerate(prs.slides):
        slide_data: dict = {
            "number":  slide_num + 1,
            "title":   "",
            "content": [],
            "images":  [],
            "notes":   "",
        }

        # BUG-FIX 15: slide.shapes.title can be None on blank/custom layouts
        try:
            title_shape = slide.shapes.title
        except Exception:
            title_shape = None

        for shape in _iter_shapes(slide.shapes):   # BUG-FIX 16: includes groups

            # ── Text extraction ───────────────────────────────────────────
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if not text:
                    continue

                # BUG-FIX 15: safe title comparison — never raise AttributeError
                if title_shape is not None and shape is title_shape:
                    slide_data["title"] = text
                else:
                    # Preserve paragraph-level text with run formatting hints
                    paragraphs = []
                    for para in shape.text_frame.paragraphs:
                        para_text = para.text.strip()
                        if para_text:
                            paragraphs.append(para_text)

                    slide_data["content"].append({
                        "type":       "text",
                        "content":    text,
                        "paragraphs": paragraphs,
                    })

            # ── Image extraction ──────────────────────────────────────────
            # BUG-FIX 16: also catch LINKED_PICTURE (type 3) alongside PICTURE (13)
            if shape.shape_type in (
                MSO_SHAPE_TYPE.PICTURE,          # 13
                MSO_SHAPE_TYPE.LINKED_PICTURE,   # 3
            ):
                try:
                    image      = shape.image
                    image_ext  = image.ext or "png"
                    image_name = (
                        f"slide{slide_num + 1}"
                        f"_img{len(slide_data['images']) + 1}"
                        f".{image_ext}"
                    )
                    image_path = os.path.join(assets_dir, image_name)

                    with open(image_path, "wb") as fh:
                        fh.write(image.blob)

                    slide_data["images"].append({
                        "path":   f"assets/{image_name}",
                        "width":  shape.width,
                        "height": shape.height,
                    })
                except Exception as img_err:
                    # Linked images with broken paths raise errors — skip silently
                    print(f"  Slide {slide_num + 1}: skipped image ({img_err})")

        # ── Speaker notes ─────────────────────────────────────────────────
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            notes_text  = notes_frame.text.strip()

            # BUG-FIX 17: PowerPoint inserts "(Click to add notes)" when empty
            PPTX_PLACEHOLDER = "(click to add notes)"
            if (notes_text
                    and notes_text.lower() != PPTX_PLACEHOLDER
                    and notes_text.lower() != "click to add notes"):
                slide_data["notes"] = notes_text

        slides_data.append(slide_data)

    return slides_data


# ─── Entry point ─────────────────────────────────────────────────────────────

def main() -> None:
    if len(sys.argv) < 2:
        print("Usage: python extract-pptx.py <input.pptx> [output_dir]")
        sys.exit(1)

    input_file = sys.argv[1]
    output_dir = sys.argv[2] if len(sys.argv) > 2 else "."

    if not os.path.exists(input_file):
        print(f"File not found: {input_file}")
        sys.exit(1)

    _ensure_pptx()   # BUG-FIX 19: auto-install

    os.makedirs(output_dir, exist_ok=True)
    slides = extract_pptx(input_file, output_dir)

    output_path = os.path.join(output_dir, "extracted-slides.json")
    with open(output_path, "w", encoding="utf-8") as fh:
        json.dump(slides, fh, indent=2, ensure_ascii=False)

    print(f"Extracted {len(slides)} slides → {output_path}")
    for s in slides:
        img_count  = len(s["images"])
        note_flag  = " [notes]" if s["notes"] else ""
        print(
            f"  Slide {s['number']:02d}: "
            f"{s['title'] or '(no title)'}"
            f" — {img_count} image(s){note_flag}"
        )


if __name__ == "__main__":
    main()
